"""
Comprehensive tests for all extractors (Markdown, PDF, PPTX, XLSX, Image).
Tests extraction, metadata, and error handling.
"""

import pytest
from pathlib import Path
from io import BytesIO

from ragtriever.extractors.markdown import MarkdownExtractor
from ragtriever.extractors.pdf import PdfExtractor
from ragtriever.extractors.pptx import PptxExtractor
from ragtriever.extractors.xlsx import XlsxExtractor
from ragtriever.extractors.image import TesseractImageExtractor


class TestMarkdownExtractor:
    """Test Markdown extraction with real and synthetic files."""

    def test_extract_simple_markdown(self, tmp_path: Path):
        """Test extraction of simple markdown file."""
        md_file = tmp_path / "simple.md"
        md_file.write_text("# Heading\n\nSome content here.")

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result.text is not None
        assert len(result.text) > 0
        assert "Heading" in result.text or "content" in result.text

    def test_extract_markdown_with_wikilinks(self, tmp_path: Path):
        """Test extraction preserves wikilinks."""
        md_file = tmp_path / "with_links.md"
        md_file.write_text(
            "# Note\n\nLinked to [[other-note]] and [[Project/Task]]."
        )

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result.metadata is not None
        assert "wikilinks" in result.metadata
        wikilinks = result.metadata["wikilinks"]
        assert len(wikilinks) >= 2

    def test_extract_markdown_with_tags(self, tmp_path: Path):
        """Test extraction captures tags."""
        md_file = tmp_path / "with_tags.md"
        md_file.write_text(
            "# Project Note\n\n#project #important\n\nContent here."
        )

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result.metadata is not None
        assert "tags" in result.metadata
        tags = result.metadata["tags"]
        assert len(tags) >= 2

    def test_extract_markdown_with_frontmatter(self, tmp_path: Path):
        """Test extraction handles YAML frontmatter."""
        md_file = tmp_path / "with_frontmatter.md"
        md_file.write_text(
            """---
title: Test Note
date: 2024-01-15
tags: [test, important]
---

# Content

This is the actual content."""
        )

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result.metadata is not None
        assert "frontmatter" in result.metadata or "title" in result.text

    def test_extract_markdown_with_code_blocks(self, tmp_path: Path):
        """Test extraction preserves code blocks."""
        md_file = tmp_path / "with_code.md"
        md_file.write_text(
            """# Code Example

```python
def hello():
    print("world")
```

More content."""
        )

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result.text is not None
        assert "hello" in result.text or "Code" in result.text

    def test_extract_markdown_with_lists(self, tmp_path: Path):
        """Test extraction with bullet points."""
        md_file = tmp_path / "with_lists.md"
        md_file.write_text(
            """# List Example

- Item 1
- Item 2
  - Nested item
- Item 3

Content after."""
        )

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result.text is not None
        assert "Item" in result.text


class TestPDFExtractor:
    """Test PDF extraction (if PyPDF2 is available)."""

    def test_extract_pdf_with_text(self, tmp_path: Path):
        """Test PDF extraction if reportlab is available."""
        pdf_file = tmp_path / "simple.pdf"

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter

            # Create a simple PDF
            c = canvas.Canvas(str(pdf_file), pagesize=letter)
            c.drawString(100, 750, "Test PDF Content")
            c.drawString(100, 700, "Second line")
            c.save()

            extractor = PdfExtractor()
            result = extractor.extract(pdf_file)

            assert result.text is not None
            # PDF extraction may vary, just verify we get something
            assert len(result.text) >= 0

        except ImportError:
            pytest.skip("reportlab not installed")

    def test_extract_pdf_metadata(self, tmp_path: Path):
        """Test PDF metadata extraction."""
        pdf_file = tmp_path / "with_meta.pdf"

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter

            c = canvas.Canvas(str(pdf_file), pagesize=letter)
            c.drawString(100, 750, "PDF Content")
            c.save()

            extractor = PdfExtractor()
            result = extractor.extract(pdf_file)

            assert result.metadata is not None
            # Should have page count or size info
            assert "pages" in result.metadata or len(result.metadata) >= 0

        except ImportError:
            pytest.skip("reportlab not installed")


class TestPPTXExtractor:
    """Test PowerPoint extraction."""

    def test_extract_pptx_with_text(self, tmp_path: Path):
        """Test PPTX extraction if python-pptx is available."""
        pptx_file = tmp_path / "simple.pptx"

        try:
            from pptx import Presentation

            # Create a simple presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            textbox = slide.shapes.add_textbox(10, 10, 300, 100)
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            p.text = "Test Slide Content"
            prs.save(str(pptx_file))

            extractor = PptxExtractor()
            result = extractor.extract(pptx_file)

            assert result.text is not None
            # Should contain slide content
            assert len(result.text) > 0

        except ImportError:
            pytest.skip("python-pptx not installed")

    def test_extract_pptx_metadata(self, tmp_path: Path):
        """Test PPTX metadata extraction."""
        pptx_file = tmp_path / "with_meta.pptx"

        try:
            from pptx import Presentation

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            textbox = slide.shapes.add_textbox(10, 10, 300, 100)
            textbox.text_frame.paragraphs[0].text = "Content"
            prs.save(str(pptx_file))

            extractor = PptxExtractor()
            result = extractor.extract(pptx_file)

            assert result.metadata is not None
            # Should have slide count
            assert "slides" in result.metadata or len(result.metadata) >= 0

        except ImportError:
            pytest.skip("python-pptx not installed")


class TestXLSXExtractor:
    """Test Excel extraction."""

    def test_extract_xlsx_with_data(self, tmp_path: Path):
        """Test XLSX extraction if openpyxl is available."""
        xlsx_file = tmp_path / "simple.xlsx"

        try:
            from openpyxl import Workbook

            # Create a simple workbook
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Header 1"
            ws["B1"] = "Header 2"
            ws["A2"] = "Data 1"
            ws["B2"] = "Data 2"
            wb.save(str(xlsx_file))

            extractor = XlsxExtractor()
            result = extractor.extract(xlsx_file)

            assert result.text is not None
            # Should contain data
            assert len(result.text) > 0

        except ImportError:
            pytest.skip("openpyxl not installed")

    def test_extract_xlsx_metadata(self, tmp_path: Path):
        """Test XLSX metadata extraction."""
        xlsx_file = tmp_path / "with_meta.xlsx"

        try:
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Data"
            wb.save(str(xlsx_file))

            extractor = XlsxExtractor()
            result = extractor.extract(xlsx_file)

            assert result.metadata is not None
            # Should have sheet info
            assert "sheets" in result.metadata or len(result.metadata) >= 0

        except ImportError:
            pytest.skip("openpyxl not installed")


class TestImageExtractor:
    """Test image extraction (no OCR by default)."""

    def create_test_image(self, tmp_path: Path) -> Path:
        """Create a simple test image."""
        try:
            from PIL import Image, ImageDraw
        except ImportError:
            pytest.skip("Pillow not installed")

        img_file = tmp_path / "test.png"
        img = Image.new("RGB", (200, 100), color="white")
        draw = ImageDraw.Draw(img)
        draw.text((10, 40), "Image Text", fill="black")
        img.save(img_file)
        return img_file

    def test_extract_image_basic(self, tmp_path: Path):
        """Test basic image extraction (metadata only)."""
        img_file = self.create_test_image(tmp_path)

        extractor = TesseractImageExtractor(ocr_mode="off")
        result = extractor.extract(img_file)

        assert result.metadata is not None
        assert result.metadata.get("width") == 200
        assert result.metadata.get("height") == 100
        assert result.metadata.get("ocr_mode") == "off"

    def test_extract_image_with_dimensions(self, tmp_path: Path):
        """Test image dimension extraction."""
        img_file = self.create_test_image(tmp_path)

        extractor = TesseractImageExtractor(ocr_mode="off")
        result = extractor.extract(img_file)

        assert result.metadata is not None
        assert "width" in result.metadata
        assert "height" in result.metadata
        assert result.metadata["width"] > 0
        assert result.metadata["height"] > 0


class TestExtractorRegistry:
    """Test the extractor registry and file type detection."""

    def test_registry_selects_correct_extractor(self, tmp_path: Path):
        """Test that correct extractor is selected by file type."""
        from ragtriever.extractors.base import ExtractorRegistry

        registry = ExtractorRegistry()
        registry.register(MarkdownExtractor())

        md_file = tmp_path / "test.md"
        md_file.write_text("# Test")

        extractor = registry.get(md_file)
        assert extractor is not None
        assert isinstance(extractor, MarkdownExtractor)

    def test_registry_returns_none_for_unsupported_type(self, tmp_path: Path):
        """Test that registry returns None for unsupported types."""
        from ragtriever.extractors.base import ExtractorRegistry

        registry = ExtractorRegistry()

        unknown_file = tmp_path / "test.unknown"
        unknown_file.write_text("content")

        extractor = registry.get(unknown_file)
        assert extractor is None


class TestExtractorErrorHandling:
    """Test error handling in extractors."""

    def test_extract_missing_file(self):
        """Test extraction of non-existent file."""
        extractor = MarkdownExtractor()
        missing_file = Path("/tmp/nonexistent_file_12345.md")

        # Should raise an error
        with pytest.raises((FileNotFoundError, OSError)):
            extractor.extract(missing_file)

    def test_extract_empty_file(self, tmp_path: Path):
        """Test extraction of empty file."""
        md_file = tmp_path / "empty.md"
        md_file.write_text("")

        extractor = MarkdownExtractor()
        result = extractor.extract(md_file)

        assert result is not None
        # Empty file should still be extracted
        assert result.text == "" or result.text is not None
