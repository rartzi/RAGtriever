from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from .base import Extracted

@dataclass
class PptxExtractor:
    supported_suffixes = (".pptx",)

    def extract(self, path: Path) -> Extracted:
        """Extract PowerPoint slide text.

        TODO: Improve structure, include speaker notes, preserve slide boundaries.
        """
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:
            raise RuntimeError("python-pptx required for PPTX extraction") from e

        prs = Presentation(str(path))
        slides_out: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"[[[SLIDE {idx}]]]"]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        parts.append(txt)
            slides_out.append("\n".join(parts))
        text = "\n\n".join(slides_out)
        meta: dict[str, Any] = {"slide_count": len(prs.slides)}
        return Extracted(text=text, metadata=meta)
