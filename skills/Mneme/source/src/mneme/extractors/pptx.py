from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterator
import logging
import zipfile
import xml.etree.ElementTree as ET

from .base import Extracted

logger = logging.getLogger(__name__)


@dataclass
class PptxExtractor:
    supported_suffixes = (".pptx",)
    min_image_width: int = 100
    min_image_height: int = 100

    def _iter_shapes(self, shapes: Any) -> Iterator[Any]:
        """Recursively iterate through all shapes, including those inside groups.

        This ensures we capture text from grouped text boxes, which are common
        in acknowledgement slides and multi-column layouts.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        except Exception:
            # If we can't import, just iterate top-level
            yield from shapes
            return

        for shape in shapes:
            yield shape
            # Recursively iterate into group shapes
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    yield from self._iter_shapes(shape.shapes)
            except (NotImplementedError, AttributeError):
                # Shape type not recognized or no shapes attribute
                pass

    def _extract_smartart_text(self, path: Path) -> dict[int, list[str]]:
        """Extract text from SmartArt diagrams in PPTX.

        SmartArt diagrams are stored as graphicFrame elements with text data
        in separate /ppt/diagrams/data*.xml files. This method:
        1. Reads slide relationship files to find diagram references
        2. Extracts text from diagram data XML files
        3. Returns a mapping of {slide_number: [diagram_texts]}
        """
        smartart_by_slide: dict[int, list[str]] = {}

        try:
            with zipfile.ZipFile(path, 'r') as zf:
                # Namespace for DrawingML
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

                # Get list of all diagram data files
                diagram_files = [name for name in zf.namelist() if name.startswith('ppt/diagrams/data') and name.endswith('.xml')]

                # Build a mapping of diagram ID to text content
                diagram_texts: dict[str, list[str]] = {}
                for diagram_file in diagram_files:
                    try:
                        xml_content = zf.read(diagram_file)
                        root = ET.fromstring(xml_content)

                        # Extract all text elements
                        texts = []
                        for t_elem in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
                            if t_elem.text:
                                text = t_elem.text.strip()
                                # Filter out placeholder text and empty strings
                                if text and not text.startswith('[') and len(text) > 1:
                                    texts.append(text)

                        if texts:
                            # Use filename as key (e.g., "data1.xml")
                            diagram_id = diagram_file.split('/')[-1]
                            diagram_texts[diagram_id] = texts

                    except Exception as e:
                        logger.warning(f"Failed to parse diagram {diagram_file}: {e}")
                        continue

                # Now map diagrams to slides via relationship files
                slide_rel_files = [name for name in zf.namelist() if name.startswith('ppt/slides/_rels/slide') and name.endswith('.xml.rels')]

                for rel_file in slide_rel_files:
                    try:
                        # Extract slide number from filename (e.g., "ppt/slides/_rels/slide2.xml.rels" -> 2)
                        filename = rel_file.split('/')[-1]  # Get "slide2.xml.rels"
                        slide_num = int(filename.replace('slide', '').replace('.xml.rels', ''))

                        # Parse relationships to find diagram references
                        xml_content = zf.read(rel_file)
                        root = ET.fromstring(xml_content)

                        # Look for relationships pointing to diagram files
                        for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                            target = rel.get('Target', '')
                            if '../diagrams/data' in target:
                                # Extract diagram filename (e.g., "../diagrams/data1.xml" -> "data1.xml")
                                diagram_id = target.split('/')[-1]
                                if diagram_id in diagram_texts:
                                    if slide_num not in smartart_by_slide:
                                        smartart_by_slide[slide_num] = []
                                    smartart_by_slide[slide_num].extend(diagram_texts[diagram_id])

                    except Exception as e:
                        logger.warning(f"Failed to parse slide relationships {rel_file}: {e}")
                        continue

        except Exception as e:
            logger.warning(f"Failed to extract SmartArt text from {path}: {e}")

        return smartart_by_slide

    def extract(self, path: Path) -> Extracted:
        """Extract PowerPoint slide text and embedded images.

        Returns text with slide markers and metadata including:
        - slide_count: number of slides
        - embedded_images: list of image metadata dicts (slide_num, image_bytes, dimensions)
        """
        try:
            from pptx import Presentation  # type: ignore
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        except Exception as e:
            raise RuntimeError("python-pptx required for PPTX extraction") from e

        prs = Presentation(str(path))
        slides_out: list[str] = []
        embedded_images: list[dict[str, Any]] = []

        # Extract SmartArt diagram text (stored separately from shapes)
        smartart_texts = self._extract_smartart_text(path)

        for idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"[[[SLIDE {idx}]]]"]

            # Use recursive shape iteration to capture grouped text boxes
            for shape in self._iter_shapes(slide.shapes):
                # Extract text
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        parts.append(txt)

                # Extract images (wrap shape_type access in try/except for unsupported shapes)
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_data = self._extract_image_from_shape(shape, slide_num=idx)
                        if img_data:
                            embedded_images.append(img_data)
                except NotImplementedError:
                    # Some shape types are not recognized by python-pptx
                    pass
                except Exception as e:
                    logger.warning(f"Failed to extract image from slide {idx}: {e}")

            # Add SmartArt diagram text for this slide
            if idx in smartart_texts:
                parts.extend(smartart_texts[idx])

            slides_out.append("\n".join(parts))

        text = "\n\n".join(slides_out)
        meta: dict[str, Any] = {
            "slide_count": len(prs.slides),
            "embedded_images": embedded_images,
        }
        return Extracted(text=text, metadata=meta)

    def _extract_image_from_shape(self, shape: Any, slide_num: int) -> dict[str, Any] | None:
        """Extract image data from a PowerPoint shape.

        Returns image metadata dict or None if image is too small:
        - slide_num: slide number
        - image_bytes: raw image data
        - width, height: dimensions
        - left, top: position on slide
        """
        try:
            # Get image blob (raw bytes)
            image_bytes = shape.image.blob

            # Get dimensions (in EMUs - English Metric Units, convert to pixels at 96 DPI)
            # 1 inch = 914400 EMUs, 96 DPI means 1 inch = 96 pixels
            # So: pixels = EMUs / 914400 * 96 = EMUs / 9525
            width_px = int(shape.width / 9525) if hasattr(shape, "width") else 0
            height_px = int(shape.height / 9525) if hasattr(shape, "height") else 0

            # Filter by size: skip tiny images
            if width_px < self.min_image_width or height_px < self.min_image_height:
                return None

            # Get position on slide
            left_px = int(shape.left / 9525) if hasattr(shape, "left") else 0
            top_px = int(shape.top / 9525) if hasattr(shape, "top") else 0

            return {
                "slide_num": slide_num,
                "image_bytes": image_bytes,
                "width": width_px,
                "height": height_px,
                "left": left_px,
                "top": top_px,
            }

        except Exception as e:
            logger.warning(f"Failed to extract image from shape on slide {slide_num}: {e}")
            return None
